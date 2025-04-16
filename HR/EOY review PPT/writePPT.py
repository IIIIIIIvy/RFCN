import pandas as pd
from tkinter import messagebox
import tkinter as tk
import os
from pathlib import Path
from datetime import datetime,timedelta
from pptx import Presentation
from pptx.util import Pt,Inches
from pptx.enum.text import PP_ALIGN

# 修改日志：
# - 2/19/2025 修改bug：由于.apply(lambda x:f"{x:,}")不改变原有数据类型，今年的bonus数据在excel中原来的格式会影响转换后的格式
#                   因此在前面多加了.astype(int)

def data_extraction(file_dir, data_file_name, year):
    data=pd.read_excel(file_dir+data_file_name,)
    for col in [str(year)+' Bonus Potential_RMB',str(year)+' Actual Bonus_RMB','Current Salary',str(year+1)+' Salary', str(year+1)+' Bonus Potential_RMB']:
        data[col]=data[col].astype(int).apply(lambda x:f"{x:,}")
    print('-----------------------------------------------------')
    print('Finish reading data.')
    return data


def write_PPT(data, folder_dir, year):
    template_file_name=str(year-1)+' EOY Template for '
    result_file_name = str(year-1)+' EOY Review for '

    print('-----------------------------------------------------')
    for dept in data['Dept.'].unique():
        # 打开该部门的模板ppt
        pre=Presentation(folder_dir+'template\\'+template_file_name+dept+'.pptx')
        # 获得该部门数据明细
        current_data=data[data['Dept.']==dept].reset_index(drop=True)
        for index,row in current_data.iterrows():
            # ----------------更改首页姓名
            pre.slides[index*7+0].shapes[0].text=row['Name']

            # ----------------插入图片
            salary_letter_name=folder_dir+"salary letter\\"+row['Name']+'.jpg'
            slide_width = pre.slide_width
            slide_height = pre.slide_height

            dpi=150
            img_width=1275/dpi
            img_height=1650/dpi
            max_width = slide_width.inches
            max_height = slide_height.inches
            if img_width > max_width or img_height > max_height:
                scale = min(max_width / img_width, max_height / img_height)
                scale*=0.8
                img_width *= scale
                img_height *= scale
            left = (slide_width - Inches(img_width)) / 2
            top = (slide_height - Inches(img_height)) / 2
            pic = pre.slides[index*7+5].shapes.add_picture(salary_letter_name, left, top, width=Inches(img_width), height=Inches(img_height))
            
            # ----------------更改数据
            for index_salary in [index*7+2, index*7+4]:
                for shape in pre.slides[index_salary].shapes:
                    if shape.has_table:
                        table=shape.table
                        
                        col_name=table.cell(0,0).text
                        col_name=col_name if col_name.find('Your')==-1 else col_name.replace('Your ','')
                        if col_name.find('Salary')!=-1:
                            table.cell(1,0).text=table.cell(1,0).text[:3]+' '+row[col_name]+' '+table.cell(1,0).text[3:]
                        elif col_name.find('Bonus Potential')!=-1:
                            bonus=row[col_name+'_RMB']
                            month=str(row[col_name+'_Month'])
                            table.cell(1,0).text=table.cell(1,0).text[:3]+' '+bonus+table.cell(1,0).text[4]+table.cell(1,0).text[4]+month+' '+table.cell(1,0).text[5:]
                        elif col_name.find('Actual Bonus')!=-1:
                            bonus=row[col_name+'_RMB']
                            month=str(row[col_name+'_Month'])
                            per=str(row[col_name+'_Per'])
                            table.cell(1,0).text=table.cell(1,0).text[:3]+' '+bonus+table.cell(1,0).text[4]+table.cell(1,0).text[4]+month+' '+table.cell(1,0).text[5:12]+per+table.cell(1,0).text[12:]
                        else:
                            table.cell(1,0).text=str(row[col_name])+' '+table.cell(1,0).text
                        
                        table.cell(1,0).text_frame.paragraphs[0].font.size=Pt(32)
                        table.cell(1,0).text_frame.paragraphs[0].font.bold=True

                        for para in table.cell(1,0).text_frame.paragraphs:
                            para.alignment = PP_ALIGN.CENTER
                            
        # 保存结果ppt
        pre.save(folder_dir+'result\\'+result_file_name+dept+'.pptx')
        print('Finish writing PPT for '+dept+'.')

    

if __name__ == "__main__":
    root = tk.Tk()
    root.wm_attributes("-topmost", 1)
    root.withdraw()
    messagebox.showinfo("提示", "Starting...")

    # folder_dir = os.getcwd()+'\\EOY review PPT\\'
    folder_dir = os.getcwd()+'\\'
    file_path_list = Path(folder_dir).glob('*.xlsx') 
    file_name_list=[x.name for x in file_path_list]

    file_name_list.sort()
    latest_file_name=file_name_list[-1]

    this_year=datetime.now().year

    data = data_extraction(file_dir=folder_dir, data_file_name=latest_file_name, year=this_year)
    write_PPT(data,folder_dir,this_year)
    messagebox.showinfo("提示", "Completed.")

